import telebot
import os
from docx import Document
from pptx import Presentation

# --- TOKEN VA API KALITLAR ---
TOKEN = "8493133123:AAG4XlRunfFMgrFKLp7yREeg-apn4jT93HI"
OPENAI_API_KEY = "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"

bot = telebot.TeleBot(TOKEN)

# --- FOYDALANUVCHI BOSHLAGANDA ---
@bot.message_handler(commands=['start'])
def start(message):
    markup = telebot.types.ReplyKeyboardMarkup(resize_keyboard=True)
    markup.row("📄 Maqola", "📘 Referat")
    markup.row("🧾 Mustaqil ish", "📊 Slayd (PowerPoint)")
    bot.send_message(
        message.chat.id,
        "Salom! Quyidagilardan birini tanlang 👇",
        reply_markup=markup
    )

# --- TANLOVNI QABUL QILISH ---
@bot.message_handler(func=lambda msg: msg.text in ["📄 Maqola", "📘 Referat", "🧾 Mustaqil ish", "📊 Slayd (PowerPoint)"])
def get_topic(message):
    user_choice = message.text
    bot.send_message(message.chat.id, f"Siz tanladingiz: {user_choice}\nEndi mavzuni kiriting ✍️")
    bot.register_next_step_handler(message, lambda m: generate_file(m, user_choice))

# --- FAYL YARATISH FUNKSIYASI ---
def generate_file(message, choice):
    topic = message.text
    bot.send_message(message.chat.id, f"Mavzu qabul qilindi: *{topic}*.\nFayl tayyorlanmoqda... ⏳", parse_mode="Markdown")

    if choice == "📊 Slayd (PowerPoint)":
        prs = Presentation()
        for i in range(10):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = f"{topic} - {i+1}-sahifa"
            content.text = f"{topic} haqida qisqacha ma'lumot (slayd {i+1})"
        filename = f"{topic}_slayd.pptx"
        prs.save(filename)
    else:
        doc = Document()
        doc.add_heading(f"{topic}", level=1)
        for i in range(10):
            doc.add_paragraph(f"{i+1}. {topic} haqida yozilgan matn qismi.")
        if choice == "📄 Maqola":
            filename = f"{topic}_maqola.docx"
        elif choice == "📘 Referat":
            filename = f"{topic}_referat.docx"
        else:
            filename = f"{topic}_mustaqil.docx"
        doc.save(filename)

    with open(filename, "rb") as file:
        bot.send_document(message.chat.id, file)
    os.remove(filename)
    bot.send_message(message.chat.id, "✅ Tayyor! Fayl yuborildi.")

print("🤖 Bot ishga tushdi...")
bot.polling(non_stop=True)
